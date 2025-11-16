from fastapi import FastAPI, Request
from fastapi.middleware.cors import CORSMiddleware
import uvicorn
from pptx import Presentation
import json
from enum import Enum
from fastapi.responses import FileResponse
import requests
from io import BytesIO
import tempfile
import os
from PIL import Image
from pptxtoimages.tools import PPTXToImageConverter

app = FastAPI(
    title="PPTX Service",
    description="FastAPI boilerplate for python-service",
    version="0.1.0",
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


@app.get("/")
async def read_root():
    return {"message": "Hello from python-service FastAPI!"}

@app.get("/list-shapes")
async def list_shapes():    
    prs = Presentation("../reference.pptx")
    
    # List all shapes from all slides
    ignored_primary_keys = ["257_2", "257_5", "257_7"]
    shapes_info = []
    for slide in prs.slides:
        for shape in slide.shapes:
            primary_key = f"{slide.slide_id}_{shape.shape_id}"
            if primary_key in ignored_primary_keys:
                continue
            
            shapes_info.append({
                "primary_key": primary_key,
                "value": "image" if shape.shape_type == 13 else shape.text,
            })    

    with open("../shapes_info.json", "w") as f:
        json.dump(shapes_info, f, indent=4)
        
    return shapes_info


class ResearchType(Enum):
    image_gen = "image_gen"
    image_given = "image_given"
    text_gen = "text_gen"
    text_given = "text_given"
    subway_gen = "subway_gen"


class Research:
    def __init__(self, prompt: str, research_type: ResearchType, max_output_tokens: int | None = None, recommended_output_tokens: int | None = None):
        self.prompt = prompt
        self.research_type = research_type
        self.max_output_tokens = max_output_tokens
        self.recommended_output_tokens = recommended_output_tokens
    
    def to_dict(self):
        return {
            "prompt": self.prompt,
            "research_type": self.research_type.value,
            "max_output_tokens": self.max_output_tokens,
            "recommended_output_tokens": self.recommended_output_tokens
        }

def shape_to_research(primary_key: str) -> Research | None:
    match primary_key:
        case "256_2":
            return Research(prompt="Generate an image", research_type=ResearchType.image_gen)
        case "256_3":
            return Research(prompt="Generate a title for a report about renting in nyc.", research_type=ResearchType.text_gen, max_output_tokens=15, recommended_output_tokens=12)
        case "256_4":
            return Research(prompt="subway_1", research_type=ResearchType.subway_gen)
        case "257_2":
            return None
        case "257_3":
            return Research(prompt="chart_1", research_type=ResearchType.image_given)
        case "257_4":
            return Research(prompt="stats_1", research_type=ResearchType.text_given)
        case "257_5":
            return None
        case "257_6":
            return Research(prompt="stats_2", research_type=ResearchType.text_given)
        case "257_7":
            return None
        case "257_8":
            return Research(prompt="Generate a text about renting in nyc", research_type=ResearchType.text_gen, max_output_tokens=30, recommended_output_tokens=26)
        case _:
            return None

@app.get("/build-heuristic")
async def build_heuristic():    
    prs = Presentation("../reference.pptx")
    
    # Build the heuristic
    heuristic = {}
    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            print(shape.name)
            # Create a unique primary key by joining slide_id and shape_id
            primary_key = f"{slide.slide_id}_{shape.shape_id}"
            heuristic[primary_key] = {
                "research": shape_to_research(primary_key).to_dict() if shape_to_research(primary_key) is not None  else None,
                "bounding_box": {
                    "left": shape.left,
                    "top": shape.top,
                    "width": shape.width,
                    "height": shape.height
                }
            }

    # Save heuristic to JSON file
    with open("../heuristic.json", "w") as f:
        json.dump(heuristic, f, indent=4)
    
    return heuristic

@app.post("/update-image")
async def update_image(request: Request):
    # takes in a primary key and the image file path to replace the image

    data = await request.json()
    referenceElementKey = data["referenceElementKey"]
    image_url = data.get("imageUrl", "../image.jpg")  # Default to ../image.jpg if not provided

    # Parse the primary key to get slide_id and shape_id
    slide_id, shape_id = referenceElementKey.split("_")
    slide_id = int(slide_id)
    shape_id = int(shape_id)

    # Load the presentation
    prs = Presentation("../reference_updated.pptx")

    # Find the slide and shape by their IDs
    target_slide = None
    for slide in prs.slides:
        if slide.slide_id == slide_id:
            target_slide = slide
            break

    if target_slide is None:
        return {"error": f"Slide with id {slide_id} not found"}

    target_shape = None
    for shape in target_slide.shapes:
        if shape.shape_id == shape_id:
            target_shape = shape
            break

    if target_shape is None:
        return {"error": f"Shape with id {shape_id} not found"}

    # Verify it's a picture shape
    if target_shape.shape_type != 13:  # MSO_SHAPE_TYPE.PICTURE
        return {"error": "Shape is not a picture"}

    # Store the position and dimensions
    left = target_shape.left
    top = target_shape.top
    width = target_shape.width
    height = target_shape.height

    # Remove the old picture shape
    sp = target_shape.element
    sp.getparent().remove(sp)

    # Handle image URL - if it starts with http, download it; otherwise treat as local path
    if image_url.startswith(('http://', 'https://')):
        # Download the image from URL
        response = requests.get(image_url)
        response.raise_for_status()

        # Save to temporary file
        with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(image_url)[1] or '.jpg') as tmp_file:
            tmp_file.write(response.content)
            temp_image_path = tmp_file.name

        image_path_to_use = temp_image_path
    else:
        # Treat as local file path
        image_path_to_use = image_url

    # Get the image dimensions to calculate aspect ratio
    img = Image.open(image_path_to_use)
    img_width, img_height = img.size
    img_aspect_ratio = img_width / img_height

    # Calculate new dimensions: keep the original height, adjust width by aspect ratio
    new_height = height
    new_width = int(new_height * img_aspect_ratio)

    # Center the image horizontally within the original bounding box
    new_left = left + (width - new_width) // 2

    # Add the new picture centered with aspect ratio preserved
    target_slide.shapes.add_picture(image_path_to_use, new_left, top, new_width, new_height)

    # Clean up temporary file if it was downloaded
    if image_url.startswith(('http://', 'https://')):
        os.unlink(temp_image_path)

    # Save the updated presentation
    output_path = "../reference_updated.pptx"
    prs.save(output_path)

    # Return the file for download
    return FileResponse(
        path=output_path,
        filename="reference_updated.pptx",
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )

@app.post("/update-shape")
async def update_shape_text(request: Request):
    # takes in a primary key and the updated text value
    
    data = await request.json()
    referenceElementKey = data["referenceElementKey"]
    content = data["content"]
    
    # Parse the primary key to get slide_id and shape_id
    slide_id, shape_id = referenceElementKey.split("_")
    slide_id = int(slide_id)
    shape_id = int(shape_id)
    
    # Load the presentation
    prs = Presentation("../reference_updated.pptx")
    
    # Find the slide and shape by their IDs
    target_slide = None
    for slide in prs.slides:
        if slide.slide_id == slide_id:
            target_slide = slide
            break
    
    if target_slide is None:
        return {"error": f"Slide with id {slide_id} not found"}
    
    target_shape = None
    for shape in target_slide.shapes:
        if shape.shape_id == shape_id:
            target_shape = shape
            break
    
    if target_shape is None:
        return {"error": f"Shape with id {shape_id} not found"}
    
    if target_shape.has_text_frame:
        text_frame = target_shape.text_frame
                
        # Clear all paragraphs except the first one
        while len(text_frame.paragraphs) > 1:
            p = text_frame.paragraphs[1]._element
            p.getparent().remove(p)

        # Update the existing run's text to preserve font properties
        paragraph = text_frame.paragraphs[0]

        if paragraph.runs:
            # Update the first run's text, preserving its font
            paragraph.runs[0].text = content
        else:
            # If no runs exist, add one
            run = paragraph.add_run()
            run.text = content
    else:
        return {"error": "Shape does not have a text frame"}
    
    # Save the updated presentation
    output_path = "../reference_updated.pptx"
    prs.save(output_path)
    
    # Return the file for download
    return FileResponse(
        path=output_path,
        filename="reference_updated.pptx",
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )

@app.get("/convert-to-png")
async def convert_to_png():
    """
    Convert the reference_updated.pptx file to PNG images.
    Returns a list of paths to the generated PNG files.
    """
    # Get the absolute path to the packages directory (parent of python-service)
    current_dir = os.path.dirname(os.path.abspath(__file__))
    packages_dir = os.path.dirname(current_dir)

    pptx_path = os.path.join(packages_dir, "reference_updated.pptx")
    output_dir = os.path.join(packages_dir, "pptx_images")

    # Check if the PPTX file exists
    if not os.path.exists(pptx_path):
        return {"error": f"PPTX file not found at {pptx_path}"}

    # Create output directory if it doesn't exist
    os.makedirs(output_dir, exist_ok=True)

    # Initialize converter with output directory and convert to images
    converter = PPTXToImageConverter(pptx_path, output_dir)
    images = converter.convert()

    # Return the list of image paths
    return {
        "message": f"Converted {len(images)} slides to PNG images",
        "image_paths": images,
        "output_directory": output_dir
    }

if __name__ == "__main__":
    # Run the app with: python main.py
    uvicorn.run("main:app", host="0.0.0.0", port=8000, reload=True)
